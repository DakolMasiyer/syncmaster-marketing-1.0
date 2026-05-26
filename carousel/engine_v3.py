import os
import json
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

class AestheticEngine:
    def __init__(self, theme_path):
        with open(theme_path, 'r') as f:
            self.theme = json.load(f)
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(self.theme['dimensions']['width'] * 9525))
        self.prs.slide_height = Emu(int(self.theme['dimensions']['height'] * 9525))
        self.blank_layout = self.prs.slide_layouts[6]

    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def px(self, n):
        return Emu(int(n * 9525))

    def pt(self, n):
        return Pt(round(n * 0.75, 2))

    def apply_style(self, run, style_token):
        scale = self.theme['typography']['scale'].get(style_token, {})
        run.font.size = Pt(scale.get('size', 18) * 0.75)
        run.font.bold = scale.get('weight', 400) >= 600
        run.font.name = self.theme['typography']['font_family']
        if scale.get('uppercase'):
            run.font.all_caps = True
        
        ls_val = scale.get('tracking', self.theme['typography']['letter_spacing'])
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(ls_val * scale.get('size', 18) * 0.75 * 100)))

    def add_text_box(self, slide, text, style_token, left, top, width, height, color_hex=None, align=PP_ALIGN.LEFT):
        color = self.hex_to_rgb(color_hex or self.theme['brand']['white'])
        tb = slide.shapes.add_textbox(self.px(left), self.px(top), self.px(width), self.px(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.color.rgb = color
        self.apply_style(run, style_token)
        return tb

    def add_background(self, slide, color_hex):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.hex_to_rgb(color_hex)

    def add_asset(self, slide, asset_path, left, top, width=None, height=None):
        if os.path.exists(asset_path):
            if width and height:
                slide.shapes.add_picture(asset_path, self.px(left), self.px(top), self.px(width), self.px(height))
            else:
                slide.shapes.add_picture(asset_path, self.px(left), self.px(top))

    def generate(self, content, output_path):
        for i, slide_data in enumerate(content['slides']):
            slide = self.prs.slides.add_slide(self.blank_layout)
            self.add_background(slide, slide_data.get('bg_color', self.theme['brand']['purple']))
            
            # Smart Placement Logic
            if slide_data.get('type') == 'opener':
                self.add_text_box(slide, slide_data['headline'], 'display-2xl', 95, 450, 890, 400, align=PP_ALIGN.CENTER)
                if 'asset' in slide_data:
                    self.add_asset(slide, slide_data['asset'], (1080-400)/2, 100, 400, 300)
            
            elif slide_data.get('type') == 'stat':
                self.add_text_box(slide, slide_data['stat'], 'display-md', 95, 300, 890, 150, color_hex=self.theme['brand']['lime'], align=PP_ALIGN.CENTER)
                self.add_text_box(slide, slide_data['headline'], 'display-lg', 95, 450, 890, 200, align=PP_ALIGN.CENTER)
            
            else:
                self.add_text_box(slide, slide_data['headline'], 'display-lg', 95, 200, 890, 200)
                if 'body' in slide_data:
                    self.add_text_box(slide, slide_data['body'], 'body-lg', 95, 450, 890, 400)
                if 'asset' in slide_data:
                    self.add_asset(slide, slide_data['asset'], 95, 800, 890, 400)
            
            # Chrome
            self.add_text_box(slide, "SYNCMASTER", 'wordmark', 95, 91, 300, 50, color_hex=self.theme['brand']['white'])
            counter_text = f"{i+1:02d} / {len(content['slides']):02d}"
            self.add_text_box(slide, counter_text, 'pill', 1080-200, 1350-130, 150, 50, color_hex=self.theme['brand']['white'], align=PP_ALIGN.RIGHT)

        self.prs.save(output_path)
        print(f"Saved: {output_path}")

if __name__ == "__main__":
    demo_content = {
        "slides": [
            {
                "type": "opener",
                "headline": "Unlock Global Music Revenue.",
                "asset": "C:/Users/infon/Documents/Claude Code/Projects/syncmaster-marketing-1.0/carousel/slide_01.png",
                "bg_color": "#5252E0"
            },
            {
                "type": "stat",
                "stat": "812%",
                "headline": "Growth in Sync Demand.",
                "bg_color": "#0A0A20"
            },
            {
                "type": "standard",
                "headline": "Your Rights Matter.",
                "body": "Don't sign away your future for a one-time fee.",
                "asset": "C:/Users/infon/Documents/Claude Code/Projects/syncmaster-marketing-1.0/carousel/slide_03.png",
                "bg_color": "#5252E0"
            }
        ]
    }
    
    engine = AestheticEngine("C:/Users/infon/Documents/Claude Code/Projects/syncmaster-marketing-1.0/carousel/styles/syncmaster.json")
    engine.generate(demo_content, "C:/Users/infon/Documents/Claude Code/Projects/syncmaster-marketing-1.0/carousel/exports/designer_agent_v3.pptx")
