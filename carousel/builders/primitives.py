"""
SyncMaster Design Agent — PPTX Primitives
All python-pptx shape/text helpers. No hardcoded values — consumes tokens only.
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn

from builders.tokens import e, pt, DM_SANS, IW, CARD_L, BORDER, MUTED


# ── Background ────────────────────────────────────────────────────────────

def solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Shape radius helper ───────────────────────────────────────────────────

def _set_radius(shp, px_radius, shorter_px):
    """Set rounded-rect corner radius. px_radius=-1 → pill (max)."""
    try:
        gd = shp.element.spPr.prstGeom.avLst.find(qn('a:gd'))
        if gd is None:
            return
        if px_radius == -1:
            gd.set('fmla', 'val 50000')
        else:
            adj = min(50000, int(px_radius / (shorter_px / 2) * 50000))
            gd.set('fmla', f'val {adj}')
    except Exception:
        pass


# ── Rectangle / rounded rect ──────────────────────────────────────────────

def add_rect(slide, left, top, w, h,
             fill=None, border=None, bpt=1.0, radius_px=0):
    if radius_px:
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            e(left), e(top), e(w), e(h))
        _set_radius(shp, radius_px, min(w, h))
    else:
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            e(left), e(top), e(w), e(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bpt)
    else:
        shp.line.fill.background()
    try:
        shp.text_frame.text = ''
    except Exception:
        pass
    return shp


# ── Oval ──────────────────────────────────────────────────────────────────

def add_oval(slide, left, top, w, h, fill=None, border=None, bpt=1.0):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, e(left), e(top), e(w), e(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bpt)
    else:
        shp.line.fill.background()
    return shp


# ── Text box (single run) ─────────────────────────────────────────────────

def tb(slide, text, left, top, w, h, size, color,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True, font=None):
    box = slide.shapes.add_textbox(e(left), e(top), e(w), e(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font or DM_SANS
    run.font.size = pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


# ── Text box (multi-run, inline colour changes) ───────────────────────────

def tb2(slide, runs, left, top, w, h, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(e(left), e(top), e(w), e(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    para = tf.paragraphs[0]
    para.alignment = align
    for rd in runs:
        run = para.add_run()
        run.text = rd['t']
        run.font.name = rd.get('f', DM_SANS)
        run.font.size = pt(rd['s'])
        run.font.color.rgb = rd['c']
        run.font.bold = rd.get('b', False)
        run.font.italic = rd.get('i', False)
    return box


# ── Pill (rounded label) ──────────────────────────────────────────────────

def pill(slide, left, top, w, h, text, border_c, text_c,
         size=11, fill=None, font=None):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        e(left), e(top), e(w), e(h))
    _set_radius(shp, -1, min(w, h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border_c:
        shp.line.color.rgb = border_c
        shp.line.width = Pt(1.0)
    else:
        shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = font or DM_SANS
    run.font.size = pt(size)
    run.font.color.rgb = text_c
    run.font.bold = False
    return shp


# ── 1px horizontal divider ────────────────────────────────────────────────

def divider(slide, y, left=None, width=None, color=None):
    left  = left  if left  is not None else CARD_L
    width = width if width is not None else IW
    color = color if color is not None else BORDER
    add_rect(slide, left, y, width, 1, fill=color)


# ── Meta key-value row (2-column grid) ───────────────────────────────────

def meta_row(slide, y, left_label, left_value, right_label, right_value,
             left_accent=False):
    """
    Renders one row of the 2-col meta grid.
    left_accent=True colours the left value with ACCENT green.
    """
    from builders.tokens import IX, C2X, CW2, MUTED, TEXT, ACCENT, DM_SANS

    lbl_h, val_h = 16, 26
    val_y = y + lbl_h + 6

    # Left column
    tb(slide, left_label.upper(), IX, y, CW2, lbl_h,
       10, MUTED, font=DM_SANS)
    tb(slide, left_value, IX, val_y, CW2, val_h,
       14, ACCENT if left_accent else TEXT, font=DM_SANS)

    # Right column
    if right_label:
        tb(slide, right_label.upper(), C2X, y, CW2, lbl_h,
           10, MUTED, font=DM_SANS)
    if right_value:
        tb(slide, right_value, C2X, val_y, CW2, val_h,
           14, TEXT, font=DM_SANS)


# ── Genre tag (small dark pill) ───────────────────────────────────────────

def genre_tag(slide, text, left, top):
    """Returns the right edge x so caller can chain tags."""
    char_w  = 7          # estimated px per char at 11px DM Sans
    pad_x   = 16
    tag_w   = len(text) * char_w + pad_x * 2
    tag_h   = 32
    pill(slide, left, top, tag_w, tag_h, text, BORDER, MUTED, size=11)
    return left + tag_w  # right edge


# ── Slot indicator (dots + labels) ───────────────────────────────────────

def slot_dots(slide, total, filled, x, y, label, deadline_text):
    from builders.tokens import ACCENT, SURFACE, BORDER, MUTED, TEXT
    from builders.tokens import AMBER, AMBER_BG

    dot_size = 12
    dot_gap  = 8
    dot_cy   = y + dot_size // 2

    # Draw dots
    for i in range(total):
        dx = x + i * (dot_size + dot_gap)
        if i < filled:
            add_oval(slide, dx, y, dot_size, dot_size, fill=ACCENT)
        else:
            add_oval(slide, dx, y, dot_size, dot_size,
                     fill=SURFACE, border=BORDER, bpt=1.0)

    # Slot label
    dots_end_x = x + total * (dot_size + dot_gap) - dot_gap + 12
    tb(slide, label, dots_end_x, y - 1, 260, dot_size + 4, 11, MUTED)

    # Deadline pill — right-aligned to content right edge
    from builders.tokens import IX, IW
    pill_w = len(deadline_text) * 7 + 28
    pill_x = IX + IW - pill_w
    pill(slide, pill_x, y - 4, pill_w, dot_size + 8,
         deadline_text, None, AMBER, size=11, fill=AMBER_BG)


# ── CTA strip ─────────────────────────────────────────────────────────────

def cta_strip(slide, y, headline, subtext, btn_label):
    from builders.tokens import IX, IW, TEXT, MUTED, DM_SERIF, DM_SANS
    from builders.tokens import SURFACE, BORDER

    # Headline — DM Serif Display
    tb(slide, headline, IX, y, IW, 70, 18, TEXT, font=DM_SERIF)

    # Subtext
    tb(slide, subtext, IX, y + 78, IW, 44, 12, MUTED, wrap=True)

    # Button
    btn_w, btn_h = 220, 48
    btn_y = y + 78 + 48 + 12
    add_rect(slide, IX, btn_y, btn_w, btn_h,
             fill=SURFACE, border=BORDER, bpt=1.0, radius_px=8)
    tb(slide, btn_label, IX, btn_y + 12, btn_w, 28,
       13, TEXT, bold=True, align=PP_ALIGN.CENTER)

    return btn_y + btn_h   # returns content bottom y


# ── Canvas branding (footer) ──────────────────────────────────────────────

def branding_footer(slide):
    from builders.tokens import ACCENT, TEXT, MUTED, FOOTER_Y, DM_SANS, W

    # Green dot
    add_oval(slide, CARD_L, FOOTER_Y + 5, 10, 10, fill=ACCENT)

    # SyncMaster wordmark
    tb(slide, 'SyncMaster', CARD_L + 18, FOOTER_Y, 160, 22, 13, TEXT,
       bold=True, font=DM_SANS)

    # URL right-aligned
    url_w = 110
    tb(slide, 'syncmaster.io', W - CARD_L - url_w, FOOTER_Y, url_w, 22,
       12, MUTED, align=PP_ALIGN.RIGHT, font=DM_SANS)
