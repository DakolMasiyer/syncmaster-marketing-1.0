"""
Export all 9 SyncMaster artboards to a PPTX.
Each slide = pixel-perfect screenshot of the rendered design.
"""

import asyncio, os, io, threading, http.server, functools, zipfile
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu

PROJECT_DIR = r"C:\Users\infon\Downloads\Syncmaster-carousel-handoff-unzip\syncmaster-carousel"
SERVE_DIR   = os.path.join(PROJECT_DIR, "project")
HTML_NAME   = "SyncMaster Carousel 01.html"
PORT        = 18765

# 1080x1350px at 96dpi -> EMU
SLIDE_W = Emu(int(1080 / 96 * 914400))
SLIDE_H = Emu(int(1350 / 96 * 914400))


def start_server():
    handler = functools.partial(http.server.SimpleHTTPRequestHandler,
                                directory=SERVE_DIR)
    httpd = http.server.HTTPServer(("127.0.0.1", PORT), handler)
    threading.Thread(target=httpd.serve_forever, daemon=True).start()
    return httpd


async def screenshot_slides():
    url = f"http://127.0.0.1:{PORT}/{HTML_NAME.replace(' ', '%20')}"
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": 1600, "height": 1000})
        await page.goto(url, wait_until="networkidle")

        # Wait for all 9 artboards + fonts to settle
        await page.wait_for_function(
            "document.querySelectorAll('[data-dc-slot]').length >= 9",
            timeout=30000
        )
        await page.wait_for_timeout(3000)

        # Flatten canvas zoom transform so cards render at natural 1080x1350
        await page.evaluate("""() => {
            const world = document.querySelector('.design-canvas > div');
            if (world) { world.style.transform = 'none'; world.style.padding = '0'; }
            document.querySelectorAll('.twk-panel').forEach(el => el.style.display = 'none');
        }""")

        cards = await page.query_selector_all(".dc-card")
        pngs = []
        for i, card in enumerate(cards, 1):
            png = await card.screenshot()   # returns bytes
            pngs.append(png)
            print(f"  screenshot {i:02d}  {len(png)//1024} KB")

        await browser.close()
    return pngs


def build_pptx(pngs):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for i, png in enumerate(pngs, 1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(io.BytesIO(png), 0, 0, SLIDE_W, SLIDE_H)

    out = os.path.join(PROJECT_DIR, "SyncMaster_Carousel_01.pptx")
    prs.save(out)
    return out


async def main():
    httpd = start_server()
    print("Rendering slides...")
    pngs = await screenshot_slides()
    httpd.shutdown()

    print("Building PPTX...")
    out = build_pptx(pngs)
    kb = os.path.getsize(out) // 1024
    print(f"\nDone -> SyncMaster_Carousel_01.pptx  ({kb} KB)")
    print(f"Location: {PROJECT_DIR}")


asyncio.run(main())
