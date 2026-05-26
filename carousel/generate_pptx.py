"""
SyncMaster Carousel 01 — Design-Token PPTX Generator
------------------------------------------------------
Every element is a real PPTX component (text box, shape, rounded rect).
Nothing is a screenshot. All text is editable in Canva.
Tokens come from chrome.jsx, slides.jsx, and linkedist-carousel-audit.json.

Unit contract
  px(n)  : CSS pixels -> EMU  (slide is 1080 px wide at 96 dpi)
  sp(n)  : CSS px font size -> PPTX Pt  (ratio = 810pt / 1080px = 0.75)
  ls(em, px_size) : em letter-spacing -> hundredths-of-pt (OOXML spc attr)
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
from PIL import Image, ImageDraw
import io, os

# ── Unit helpers ──────────────────────────────────────────────────────────────

def px(n):
    """CSS pixels -> EMU (1 px = 9525 EMU at 96 dpi)."""
    return Emu(int(n * 9525))

def sp(n):
    """CSS px font size -> PPTX Pt. Slide is 1080 px = 810 pt wide."""
    return Pt(round(n * 0.75, 2))

def ls(em_val, px_size):
    """Letter-spacing in hundredths-of-pt for OOXML spc attr."""
    return int(em_val * px_size * 0.75 * 100)

# ── Slide dimensions ──────────────────────────────────────────────────────────

W = px(1080)
H = px(1350)

# ── Design tokens — colours ───────────────────────────────────────────────────

def rgb(r, g, b): return RGBColor(r, g, b)

C_BG      = rgb(0xF9, 0xF9, 0xF9)   # warm off-white
C_INK     = rgb(0x00, 0x00, 0x00)
C_ACCENT  = rgb(0x36, 0x61, 0xFE)   # Linkedist blue
C_PILL    = rgb(0x0A, 0x0A, 0x0A)   # pill stroke / fill
C_CARD    = rgb(0xFF, 0xFF, 0xFF)
C_BORDER  = rgb(0xE0, 0xE0, 0xE0)
C_MUTED   = rgb(0x80, 0x80, 0x80)
C_DIMMED  = rgb(0x60, 0x60, 0x60)
C_WHITE   = rgb(0xFF, 0xFF, 0xFF)
C_PLACEHOLDER = rgb(0xCC, 0xCC, 0xCC)

# ── Design tokens — typography ────────────────────────────────────────────────
# Each dict: pt (Pt obj), bold, ls_spc (hundredths-of-pt), lh (thousandths-%), caps

DISPLAY_XL = dict(pt=sp(96), bold=True,  spc=ls(-0.025, 96), lh=102000, caps=False)
DISPLAY_LG = dict(pt=sp(84), bold=True,  spc=ls(-0.025, 84), lh=102000, caps=False)
DISPLAY_MD = dict(pt=sp(64), bold=False, spc=ls(-0.025, 64), lh=102000, caps=False)
EYEBROW    = dict(pt=sp(22), bold=False, spc=ls( 0.18,  22), lh=None,   caps=True)
WORDMARK   = dict(pt=sp(22), bold=True,  spc=ls( 0.10,  22), lh=None,   caps=True)
PILL_T     = dict(pt=sp(22), bold=False, spc=ls( 0.04,  22), lh=None,   caps=True)
CAPTION    = dict(pt=sp(28), bold=False, spc=0,              lh=135000, caps=False)
ROW_LABEL  = dict(pt=sp(28), bold=False, spc=ls( 0.08,  28), lh=None,   caps=True)
ROW_VAL    = dict(pt=sp(72), bold=True,  spc=ls(-0.02,  72), lh=None,   caps=False)
BUYOUT     = dict(pt=sp(124),bold=True,  spc=ls(-0.04, 124), lh=None,   caps=False)
BUYOUT_SUB = dict(pt=sp(26), bold=False, spc=ls( 0.06,  26), lh=None,   caps=True)
CARD_TITLE = dict(pt=sp(30), bold=True,  spc=ls(-0.02,  30), lh=None,   caps=False)
CARD_CAP   = dict(pt=sp(20), bold=False, spc=0,              lh=135000, caps=False)
CTA_BTN    = dict(pt=sp(44), bold=True,  spc=ls( 0.04,  44), lh=None,   caps=True)
CTA_SUB    = dict(pt=sp(36), bold=False, spc=0,              lh=None,   caps=False)
COUNTER_XL = dict(pt=sp(220),bold=True,  spc=0,              lh=None,   caps=False)

# ── Layout constants (from chrome.jsx) ───────────────────────────────────────

STAGE_TOP = px(150)
STAGE_BOT = H - px(150)
STAGE_MID = px(675)       # vertical centre of stage  (150 + 1050/2)
MARGIN    = px(95)
STAGE_W   = W - 2 * MARGIN


# ── Dot-grid background (raster — CSS radial-gradient has no PPTX equivalent) ──

def _dot_grid_png():
    img = Image.new("RGB", (1080, 1350), (0xF9, 0xF9, 0xF9))
    d = ImageDraw.Draw(img)
    for y in range(0, 1350, 35):
        for x in range(0, 1080, 35):
            d.ellipse([x - 1, y - 1, x + 2, y + 2], fill=(0xE6, 0xE6, 0xE6))
    buf = io.BytesIO()
    img.save(buf, "PNG", optimize=True)
    buf.seek(0)
    return buf

_GRID = _dot_grid_png()


# ── Low-level PPTX helpers ────────────────────────────────────────────────────

def _send_back(slide, el):
    t = slide.shapes._spTree
    t.remove(el)
    t.insert(2, el)


def _apply_run(run, tok, color=None):
    run.font.size   = tok['pt']
    run.font.bold   = tok.get('bold', False)
    run.font.name   = "Space Grotesk"
    run.font.color.rgb = color or C_INK
    if tok.get('caps'):
        run.font.all_caps = True
    spc = tok.get('spc', 0)
    if spc:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(spc))


def _apply_lh(para, tok):
    lh = tok.get('lh')
    if not lh:
        return
    pPr = para._p.get_or_add_pPr()
    for el in pPr.findall(qn('a:lnSpc')):
        pPr.remove(el)
    lnSpc  = etree.SubElement(pPr, qn('a:lnSpc'))
    spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
    spcPct.set('val', str(lh))


def _zero_para_spacing(para):
    pPr = para._p.get_or_add_pPr()
    for tag in ('a:spcBef', 'a:spcAft'):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    for tag, attr in [('a:spcBef', 'a:spcPts'), ('a:spcAft', 'a:spcPts')]:
        parent = etree.SubElement(pPr, qn(tag))
        child  = etree.SubElement(parent, qn(attr))
        child.set('val', '0')


def _rounded(slide, left, top, w, h, radius_px,
             fill=None, border=None, border_px=2):
    shape = slide.shapes.add_shape(1, left, top, w, h)
    sp    = shape._element
    spPr  = sp.find(qn('p:spPr'))
    for tag in (qn('a:prstGeom'), qn('a:custGeom')):
        el = spPr.find(tag)
        if el is not None:
            spPr.remove(el)
    min_dim = min(w / 9525, h / 9525)
    adj = min(50000, int(radius_px / min_dim * 100000)) if radius_px < min_dim else 50000
    spPr.insert(0, parse_xml(
        f'<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'prst="roundRect"><a:avLst>'
        f'<a:gd name="adj" fmla="val {adj}"/>'
        f'</a:avLst></a:prstGeom>'
    ))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Emu(border_px * 9525)
    else:
        shape.line.fill.background()
    return shape


# ── High-level component builders ─────────────────────────────────────────────

def add_bg(slide):
    _GRID.seek(0)
    pic = slide.shapes.add_picture(_GRID, 0, 0, W, H)
    _send_back(slide, pic._element)


def text_box(slide, runs_by_line, left, top, width, height,
             align=PP_ALIGN.LEFT, wrap=True):
    """
    runs_by_line: list of lines.
    Each line is either:
      - (text, tok, color)            single run
      - [(text, tok, color), ...]     multiple runs on one paragraph
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

    for i, line in enumerate(runs_by_line):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        items = line if isinstance(line, list) else [line]
        _apply_lh(p, items[0][1])
        _zero_para_spacing(p)
        for text, tok, color in items:
            r = p.add_run()
            r.text = text
            _apply_run(r, tok, color)
    return tb


def pill(slide, text, left, top, width, height, filled=False):
    shape = _rounded(slide, left, top, width, height, radius_px=9999,
                     fill=C_PILL if filled else None,
                     border=C_PILL, border_px=2)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.margin_left = tf.margin_right = px(20)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _apply_run(r, PILL_T, C_WHITE if filled else C_INK)
    return shape


def card(slide, left, top, width, height, radius_px=40):
    return _rounded(slide, left, top, width, height, radius_px,
                    fill=C_CARD, border=C_BORDER, border_px=1)


def image_placeholder(slide, label, left, top, width, height, radius_px=40):
    _rounded(slide, left, top, width, height, radius_px,
             fill=C_PLACEHOLDER, border=rgb(0xAA, 0xAA, 0xAA), border_px=1)
    tok = dict(pt=sp(20), bold=False, spc=0, lh=None, caps=False)
    text_box(slide, [(f"[ {label} ]", tok, C_MUTED)],
             left, top + (height - px(30)) // 2, width, px(30),
             align=PP_ALIGN.CENTER)


def chrome(slide, num, category, show_swipe=True, center_wm=False):
    pill_h = px(46)
    pill_y = H - px(79) - pill_h

    # Wordmark
    wm_w = px(210)
    wm_x = (W - wm_w) // 2 if center_wm else MARGIN
    text_box(slide, [("SYNCMASTER", WORDMARK, C_INK)],
             wm_x, px(84), wm_w, px(36))

    # Category pill (bottom-left)
    cat_w = px(min(320, max(180, len(category) * 13 + 56)))
    pill(slide, category, MARGIN, pill_y, cat_w, pill_h)

    # Counter pill (bottom-right)
    cnt_w = px(118)
    pill(slide, f"{num:02d} / 09", W - MARGIN - cnt_w, pill_y, cnt_w, pill_h)

    # Swipe arrow (top-right)
    if show_swipe:
        sw_w, sw_h = px(88), px(46)
        shape = _rounded(slide, W - MARGIN - sw_w, px(79), sw_w, sw_h,
                         radius_px=9999, border=C_PILL, border_px=2)
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "→"          # →
        r.font.size = sp(22)
        r.font.name = "Space Grotesk"
        r.font.color.rgb = C_INK


def divider(slide, left, top, width):
    shape = slide.shapes.add_shape(1, left, top, width, px(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_BORDER
    shape.line.fill.background()


# ── Slide builders ────────────────────────────────────────────────────────────

def slide1(slide):
    """Hook — You scored 4 Nollywood films."""
    add_bg(slide)
    chrome(slide, 1, "THE COMPOSER PROBLEM")

    hl_top = STAGE_TOP + px(30)
    text_box(slide, [
        [("You scored ", DISPLAY_XL, C_INK),
         ("4 Nollywood films.", DISPLAY_XL, C_ACCENT)],
        [("None of them pay you again.", DISPLAY_XL, C_INK)],
    ], MARGIN, hl_top, STAGE_W, px(310))

    image_placeholder(slide, "Greyscale Lagos / film-set photo",
                      MARGIN, hl_top + px(330), STAGE_W, px(360))


def slide2(slide):
    """Walk-in — A producer called you on a Tuesday."""
    add_bg(slide)
    chrome(slide, 2, "THE COMPOSER PROBLEM")

    eyebrow_y = STAGE_MID - px(120)
    ew = px(260)
    text_box(slide, [("SCENE ONE", EYEBROW, C_MUTED)],
             (W - ew) // 2, eyebrow_y, ew, px(40), align=PP_ALIGN.CENTER)

    text_box(slide, [("A producer called you on a Tuesday.", DISPLAY_XL, C_INK)],
             MARGIN, eyebrow_y + px(60), STAGE_W, px(300), align=PP_ALIGN.CENTER)


def slide3(slide):
    """The Deal — 12 cues / 3 weeks / 180,000 full buyout."""
    add_bg(slide)
    chrome(slide, 3, "THE COMPOSER PROBLEM")

    card_w = px(740)
    card_h = px(500)
    card_x = (W - card_w) // 2
    card_y = STAGE_MID - card_h // 2 - px(20)

    card(slide, card_x, card_y, card_w, card_h, radius_px=40)

    cp = px(64)   # card padding
    row_h = px(56)
    row_y = card_y + cp

    def data_row(label, value, y):
        text_box(slide, [(label, ROW_LABEL, C_MUTED)],
                 card_x + cp, y, px(200), row_h)
        text_box(slide, [(value, ROW_VAL, C_INK)],
                 card_x + card_w - cp - px(380), y, px(380), row_h,
                 align=PP_ALIGN.RIGHT)
        divider(slide, card_x + cp, y + row_h + px(14), card_w - 2 * cp)

    data_row("Cues",     "12",      row_y)
    data_row("Timeline", "3 weeks", row_y + row_h + px(16) + px(2) + px(20))

    fee_y = row_y + (row_h + px(16) + px(2) + px(20)) * 2 + px(12)
    text_box(slide, [("Fee", ROW_LABEL, C_MUTED)],
             card_x + cp, fee_y, px(200), px(40))
    text_box(slide, [("₦180,000", BUYOUT, C_INK)],
             card_x, fee_y + px(44), card_w, px(120), align=PP_ALIGN.CENTER)
    text_box(slide, [("— full buyout —", BUYOUT_SUB, C_MUTED)],
             card_x, fee_y + px(164), card_w, px(40), align=PP_ALIGN.CENTER)

    cap_y = card_y + card_h + px(36)
    text_box(slide, [("Signed on WhatsApp. No contract you kept.", CAPTION, C_MUTED)],
             MARGIN, cap_y, STAGE_W, px(56), align=PP_ALIGN.CENTER)


def slide4(slide):
    """Release — The film hit Netflix."""
    add_bg(slide)
    chrome(slide, 4, "THE COMPOSER PROBLEM")

    hl_top = STAGE_TOP + px(24)
    text_box(slide, [
        [("The film hit Netflix.", DISPLAY_LG, C_INK)],
        [("Then YouTube.", DISPLAY_LG, C_INK)],
        [("Then a hotel lobby in Accra.", DISPLAY_LG, C_INK)],
    ], MARGIN, hl_top, STAGE_W, px(380))

    image_placeholder(slide, "Greyscale street / cinema / TV photo",
                      MARGIN, hl_top + px(390), STAGE_W, px(320))


def slide5(slide):
    """Silence — Your phone never rang again."""
    add_bg(slide)
    chrome(slide, 5, "THE COMPOSER PROBLEM")

    text_box(slide, [("Your phone never rang again.", DISPLAY_XL, C_INK)],
             MARGIN, STAGE_MID - px(100), STAGE_W, px(200),
             align=PP_ALIGN.CENTER)


def slide6(slide):
    """Reframe — This isn't a talent problem. It's a rights problem."""
    add_bg(slide)
    chrome(slide, 6, "THE TURN")

    text_box(slide, [
        [("This isn’t a talent problem.", DISPLAY_XL, C_INK)],
    ], MARGIN, STAGE_MID - px(180), STAGE_W, px(160), align=PP_ALIGN.CENTER)

    text_box(slide, [
        [("It’s a ", DISPLAY_XL, C_INK),
         ("rights", DISPLAY_XL, C_ACCENT),
         (" problem.", DISPLAY_XL, C_INK)],
    ], MARGIN, STAGE_MID + px(10), STAGE_W, px(160), align=PP_ALIGN.CENTER)


def slide7(slide):
    """Mechanism — 2x2 grid."""
    add_bg(slide)
    chrome(slide, 7, "WHAT WE DO")

    text_box(slide, [("What SyncMaster does:", DISPLAY_MD, C_INK)],
             MARGIN, STAGE_TOP + px(16), STAGE_W, px(100))

    gap    = px(24)
    grid_t = STAGE_TOP + px(136)
    grid_b = STAGE_BOT - px(8)
    grid_h = grid_b - grid_t
    cell_w = (STAGE_W - gap) // 2
    cell_h = (grid_h - gap) // 2

    cards_data = [
        ("Structured briefs",
         "Producers post real opportunities. You see them first."),
        ("Rights kept clean",
         "Sync licences, not buyouts. You keep the publishing."),
        ("Curated, not flooded",
         "5 composers submit per brief. Not 500."),
        ("Paid on placement",
         "When your cue is used, you get paid. Every time."),
    ]
    positions = [
        (MARGIN, grid_t),
        (MARGIN + cell_w + gap, grid_t),
        (MARGIN, grid_t + cell_h + gap),
        (MARGIN + cell_w + gap, grid_t + cell_h + gap),
    ]

    for (title, caption), (cx, cy) in zip(cards_data, positions):
        card(slide, cx, cy, cell_w, cell_h, radius_px=32)
        pad = px(36)
        # Accent line (top of card)
        accent_bar = _rounded(slide, cx + pad, cy + pad, px(32), px(4),
                              radius_px=9999, fill=C_ACCENT)
        text_box(slide, [(title, CARD_TITLE, C_INK)],
                 cx + pad, cy + px(56), cell_w - 2 * pad, px(80))
        text_box(slide, [(caption, CARD_CAP, C_DIMMED)],
                 cx + pad, cy + px(148), cell_w - 2 * pad,
                 cell_h - px(190), wrap=True)


def slide8(slide):
    """Promise — You scored it once. Now it should pay you forever."""
    add_bg(slide)
    chrome(slide, 8, "THE TURN")

    # Oversized closing quote mark
    quote_w = px(200)
    text_box(slide, [("”", COUNTER_XL, C_ACCENT)],
             (W - quote_w) // 2, STAGE_MID - px(280),
             quote_w, px(200), align=PP_ALIGN.CENTER)

    text_box(slide, [("You scored it once.", DISPLAY_XL, C_INK)],
             MARGIN, STAGE_MID - px(30), STAGE_W, px(140), align=PP_ALIGN.CENTER)

    text_box(slide, [("Now it should pay you forever.", DISPLAY_XL, C_ACCENT)],
             MARGIN, STAGE_MID + px(125), STAGE_W, px(140), align=PP_ALIGN.CENTER)


def slide9(slide):
    """CTA — Join the Waitlist."""
    add_bg(slide)
    chrome(slide, 9, "JOIN NOW", show_swipe=False, center_wm=True)

    text_box(slide, [("SyncMaster is live.", DISPLAY_XL, C_INK)],
             MARGIN, STAGE_MID - px(270), STAGE_W, px(160), align=PP_ALIGN.CENTER)

    text_box(slide, [("We’re onboarding Nigerian composers right now.",
                      CTA_SUB, C_MUTED)],
             MARGIN, STAGE_MID - px(90), STAGE_W, px(100), align=PP_ALIGN.CENTER)

    # CTA button
    btn_w, btn_h = px(580), px(100)
    btn_x = (W - btn_w) // 2
    btn_y = STAGE_MID + px(40)
    shape = _rounded(slide, btn_x, btn_y, btn_w, btn_h, radius_px=9999,
                     fill=C_ACCENT)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "JOIN THE WAITLIST  →"
    _apply_run(r, CTA_BTN, C_WHITE)

    text_box(slide, [("Link in bio. Takes 60 seconds.", CAPTION, C_MUTED)],
             MARGIN, btn_y + btn_h + px(48), STAGE_W, px(56),
             align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

SLIDES = [slide1, slide2, slide3, slide4, slide5,
          slide6, slide7, slide8, slide9]

def generate():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    for i, builder in enumerate(SLIDES, 1):
        s = prs.slides.add_slide(blank)
        builder(s)
        print(f"  slide {i:02d} done")

    out = os.path.join(os.path.dirname(__file__), "SyncMaster_Carousel_01.pptx")
    prs.save(out)
    kb = os.path.getsize(out) // 1024
    print(f"\nSaved -> SyncMaster_Carousel_01.pptx  ({kb} KB)")

if __name__ == "__main__":
    print("Building SyncMaster Carousel 01 (design-token PPTX)...")
    generate()
