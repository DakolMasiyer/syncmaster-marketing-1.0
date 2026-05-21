"""
SyncMaster Carousel 02 — Fully Editable PPTX
All elements are native pptx shapes/text frames. No screenshots.
Canva-ready: every text box, background, and pill is independently editable.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
import os

# ── px→EMU and px→Pt converters (96 dpi baseline) ────────────────────────
def e(n): return int(n * 9525)
def p(n): return Pt(n * 0.75)

W, H = 1080, 1350
FONT = "Space Grotesk"

# ── Brand colours ─────────────────────────────────────────────────────────
PURPLE = RGBColor(0x52, 0x52, 0xE0)
LIME   = RGBColor(0xC9, 0xE8, 0x34)
DARK   = RGBColor(0x0A, 0x0A, 0x20)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x0D, 0x0D, 0x18)

# Alpha-blended colours (pre-composited against bg for pptx solid fills)
PILL_P  = RGBColor(0x9B, 0x9B, 0xF1)  # rgba(fff,.42) on purple — pill border
PILL_D  = RGBColor(0x8C, 0x8C, 0x9D)  # rgba(fff,.55) on dark   — pill border
DIM_P   = RGBColor(0xAD, 0xAD, 0xF1)  # rgba(fff,.52) on purple — eyebrow/dim
BODY_D  = RGBColor(0xC4, 0xC4, 0xCE)  # rgba(fff,.75) on dark   — body copy
DIM_D   = RGBColor(0x96, 0x96, 0xA3)  # rgba(fff,.55) on dark   — eyebrow/artist
CARD_BG = RGBColor(0x18, 0x19, 0x2A)  # rgba(fff,.09) on dark   — NowPlaying card
CARD_BD = RGBColor(0x2A, 0x2A, 0x3B)  # rgba(fff,.16) on dark   — NowPlaying border
STAT_BG = RGBColor(0x5A, 0x5A, 0xE3)  # rgba(fff,.07) on purple — stat block bg
STAT_BD = RGBColor(0x64, 0x64, 0xE4)  # rgba(fff,.11) on purple — stat block border
SRC_CLR = RGBColor(0x6E, 0x6E, 0xA8)  # rgba(fff,.38) on purple — footnote
SUB_10  = RGBColor(0xB8, 0xB8, 0xD5)  # rgba(fff,.72) on purple — slide10 sub
CAP_10  = RGBColor(0x85, 0x85, 0xAE)  # rgba(fff,.48) on purple — slide10 caption
LIME_BTN_SHD = RGBColor(0xC9, 0xE8, 0x34)  # CTA button (same lime, shadow via shape)


# ════════════════════════════════════════════════════════════════════════
# Primitive helpers
# ════════════════════════════════════════════════════════════════════════

def solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_radius(shp, px_radius, shorter_px):
    """Set rounded-rect corner radius. px_radius=-1 → pill (max)."""
    try:
        gd = shp.element.spPr.prstGeom.avLst.find(qn('a:gd'))
        if gd is None:
            return
        if px_radius == -1:
            gd.set('fmla', 'val 50000')
        else:
            adj = min(50000, int(px_radius / (shorter_px / 2) * 50000))
            gd.set('fmla', f'val {adj}')
    except Exception:
        pass


def add_rect(slide, left, top, w, h, fill=None, border=None, bpt=1.5, radius_px=0):
    if radius_px:
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                      e(left), e(top), e(w), e(h))
        _set_radius(shp, radius_px, min(w, h))
    else:
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                      e(left), e(top), e(w), e(h))
    shp.fill.solid() if fill else shp.fill.background()
    if fill:
        shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bpt)
    else:
        shp.line.fill.background()
    # Clear default text
    shp.text_frame.text = ''
    return shp


def add_oval(slide, left, top, w, h, fill=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, e(left), e(top), e(w), e(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    shp.line.fill.background()
    return shp


def tb(slide, text, left, top, w, h, size, color,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    """Add a plain text box."""
    box = slide.shapes.add_textbox(e(left), e(top), e(w), e(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = p(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def tb2(slide, runs, left, top, w, h, align=PP_ALIGN.LEFT, wrap=True):
    """Text box with multiple runs (inline colour/style changes)."""
    box = slide.shapes.add_textbox(e(left), e(top), e(w), e(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    para = tf.paragraphs[0]
    para.alignment = align
    for r in runs:
        run = para.add_run()
        run.text = r['t']
        run.font.name = FONT
        run.font.size = p(r['s'])
        run.font.color.rgb = r['c']
        run.font.bold = r.get('b', False)
        run.font.italic = r.get('i', False)
    return box


def pill(slide, left, top, w, h, text, border_c, text_c, size=22, fill=None):
    """Outlined pill with centred label text."""
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  e(left), e(top), e(w), e(h))
    _set_radius(shp, -1, min(w, h))  # pill = max radius
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    shp.line.color.rgb = border_c
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = p(size)
    run.font.color.rgb = text_c
    return shp


# ════════════════════════════════════════════════════════════════════════
# Slide chrome (wordmark + swipe pill + category + counter)
# ════════════════════════════════════════════════════════════════════════

def chrome(slide, num, total, category, swipe=True, center_wm=False, dark=False):
    pb = PILL_D if dark else PILL_P

    # Wordmark
    wm_l = (W - 280) / 2 if center_wm else 95
    wm_align = PP_ALIGN.CENTER if center_wm else PP_ALIGN.LEFT
    tb(slide, 'SYNCMASTER', wm_l, 83, 280, 34, 22, WHITE, bold=True, align=wm_align)

    # Swipe arrow pill (88×44px, top-right)
    if swipe:
        pill(slide, W - 79 - 88, 79, 88, 44, '→', pb, WHITE, size=26)

    # Category pill (bottom-left) — width scales with text
    cat = category.upper()
    cat_w = max(175, len(cat) * 13 + 52)
    pill(slide, 79, H - 79 - 44, cat_w, 44, cat, pb, WHITE)

    # Counter pill (bottom-right) — fixed width
    cnt = f'{num:02d} / {total:02d}'
    pill(slide, W - 79 - 126, H - 79 - 44, 126, 44, cnt, pb, WHITE)


# ════════════════════════════════════════════════════════════════════════
# Glow orb (approximate radial lime bloom)
# ════════════════════════════════════════════════════════════════════════

def glow(slide, top_pct, size_px, opacity):
    if size_px <= 0:
        return
    o = opacity * 0.45  # tone down for printable look
    gc = RGBColor(
        int(0xC9 * o + 0x52 * (1 - o)),
        int(0xE8 * o + 0x52 * (1 - o)),
        int(0x34 * o + 0xE0 * (1 - o)),
    )
    r = size_px / 2
    cx, cy = W / 2, H * top_pct
    add_oval(slide, cx - r, cy - r, size_px, size_px, fill=gc)


# ════════════════════════════════════════════════════════════════════════
# NowPlayingCard
# ════════════════════════════════════════════════════════════════════════

def now_playing(slide, track, artist, top=450):
    """Frosted NowPlaying pill. Returns bottom-y of the card."""
    cw, ch = 890, 96
    cl = (W - cw) / 2

    # Card background (approximated frosted glass → semi-dark rounded rect)
    add_rect(slide, cl, top, cw, ch, fill=CARD_BG, border=CARD_BD, bpt=0.75, radius_px=22)

    # Lime play-button circle
    cs = 56
    cy = top + (ch - cs) / 2
    il = cl + 32
    add_oval(slide, il, cy, cs, cs, fill=LIME)
    tb(slide, '▶', il + 14, cy + 15, cs - 14, cs - 20, 20, INK, align=PP_ALIGN.CENTER)

    # Track / artist text
    tl = il + cs + 22
    tw = cw - (cs + 22 + 32 + 90)
    tb(slide, track,  tl, top + 13, tw, 40, 30, WHITE, bold=True)
    tb(slide, artist, tl, top + 57, tw, 30, 22, DIM_D)

    # Waveform bars  [12,20,8,22,14,18,10] height in px, 3px wide
    bars = [12, 20, 8, 22, 14, 18, 10]
    bar_w, gap = 3, 4
    total_bar_w = len(bars) * bar_w + (len(bars) - 1) * gap
    bx0 = cl + cw - 32 - total_bar_w
    for i, bh in enumerate(bars):
        bx = bx0 + i * (bar_w + gap)
        by = top + (ch - 26) / 2 + (26 - bh)
        add_rect(slide, bx, by, bar_w, bh, fill=LIME, radius_px=2)

    return top + ch


# ════════════════════════════════════════════════════════════════════════
# Slide builders
# ════════════════════════════════════════════════════════════════════════

def s01(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, PURPLE)
    glow(s, 0.48, 580, 0.28)
    chrome(s, 1, 10, 'DID YOU CATCH THIS?')
    tb(s, "Your favourite songs have been on the world's biggest stages.",
       95, 380, W - 190, 460, 98, WHITE, bold=True, wrap=True)
    tb(s, 'Did you catch them?',
       95, 860, W - 190, 90, 46, LIME)


def s02(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, DARK)
    # Image placeholder note
    tb(s, '[ Slide 2: Replace dark background with NBA court / arena photo ]',
       20, 20, W - 40, 30, 14, RGBColor(0x44, 0x44, 0x66))
    chrome(s, 2, 10, 'DID YOU CATCH THIS?', dark=True)
    cb = now_playing(s, 'One Dance', 'Drake ft. Wizkid & Kyla')
    et = cb + 56
    tb(s, 'DID YOU CATCH IT AT THE…', 95, et,        W - 190, 36,  22, DIM_D)
    tb(s, 'NBA Playoffs. 2016.',       95, et + 46,   W - 190, 120, 84, WHITE, bold=True)
    tb2(s, [
        {'t': "That wasn't radio. ", 's': 34, 'c': BODY_D},
        {'t': 'That was a sync deal.', 's': 34, 'c': LIME, 'i': True},
    ], 95, et + 180, W - 190, 100)


def s03(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, DARK)
    tb(s, '[ Slide 3: Replace dark background with Top Boy still or dark London street photo ]',
       20, 20, W - 40, 30, 14, RGBColor(0x44, 0x44, 0x66))
    chrome(s, 3, 10, 'DID YOU CATCH THIS?', dark=True)
    cb = now_playing(s, 'Location', 'Burna Boy')
    et = cb + 56
    tb(s, 'DID YOU CATCH IT ON…',       95, et,       W - 190, 36,  22, DIM_D)
    tb(s, 'Netflix. Top Boy. Season 1.', 95, et + 46,  W - 190, 130, 84, WHITE, bold=True)
    tb2(s, [
        {'t': 'Lagos sound. London screen. ', 's': 34, 'c': BODY_D},
        {'t': 'That was a sync deal.', 's': 34, 'c': LIME, 'i': True},
    ], 95, et + 190, W - 190, 100)


def s04(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, DARK)
    tb(s, '[ Slide 4: Replace dark background with Spider-Man: Across the Spider-Verse still ]',
       20, 20, W - 40, 30, 14, RGBColor(0x44, 0x44, 0x66))
    chrome(s, 4, 10, 'DID YOU CATCH THIS?', dark=True)
    cb = now_playing(s, 'Free Mind', 'Tems')
    et = cb + 56
    tb(s, 'DID YOU CATCH IT IN…',                95, et,      W - 190, 36,  22, DIM_D)
    tb(s, 'Spider-Man: Across the Spider-Verse.', 95, et + 46, W - 190, 210, 84, WHITE, bold=True)
    tb2(s, [
        {'t': 'End credits. Everyone Shazamed it. ', 's': 34, 'c': BODY_D},
        {'t': 'That was a sync deal.', 's': 34, 'c': LIME, 'i': True},
    ], 95, et + 270, W - 190, 100)


def s05(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, DARK)
    tb(s, '[ Slide 5: Replace dark background with UEFA Champions League broadcast / stadium photo ]',
       20, 20, W - 40, 30, 14, RGBColor(0x44, 0x44, 0x66))
    chrome(s, 5, 10, 'DID YOU CATCH THIS?', dark=True)
    cb = now_playing(s, 'Peru', 'Fireboy DML')
    et = cb + 56
    tb(s, 'DID YOU CATCH IT AT THE…',                          95, et,      W - 190, 36,  22, DIM_D)
    tb(s, 'UEFA Champions League.\n180 countries watching.',    95, et + 46, W - 190, 210, 84, WHITE, bold=True)
    tb2(s, [
        {'t': 'One track. One deal. ', 's': 34, 'c': BODY_D},
        {'t': 'One billion impressions.', 's': 34, 'c': LIME, 'i': True},
    ], 95, et + 270, W - 190, 100)


def s06(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, DARK)
    tb(s, '[ Slide 6: Replace dark background with Black Panther: Wakanda Forever poster ]',
       20, 20, W - 40, 30, 14, RGBColor(0x44, 0x44, 0x66))
    chrome(s, 6, 10, 'DID YOU CATCH THIS?', dark=True)
    cb = now_playing(s, 'Wakanda Forever OST', 'Various African Artists')
    et = cb + 56
    tb(s, 'DID YOU CATCH IT IN…',                            95, et,      W - 190, 36,  22, DIM_D)
    tb(s, 'Black Panther: Wakanda Forever. Disney+.',         95, et + 46, W - 190, 210, 84, WHITE, bold=True)
    tb2(s, [
        {'t': 'The whole film was the brief. ', 's': 34, 'c': BODY_D},
        {'t': 'That was a sync deal.', 's': 34, 'c': LIME, 'i': True},
    ], 95, et + 270, W - 190, 100)


def s07(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, PURPLE)
    glow(s, 0.58, 420, 0.22)
    chrome(s, 7, 10, 'THE TURN')
    tb(s, 'THE QUESTION',                                             95, 320, W - 190, 36, 22, DIM_P)
    tb(s, 'So what actually happened when those songs played?',       95, 370, W - 190, 330, 84, WHITE, bold=True)
    tb(s, 'This is what a sync deal does.',                           95, 720, W - 190, 75, 38, LIME)


def s08(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, PURPLE)
    glow(s, 0.10, 360, 0.20)
    chrome(s, 8, 10, 'SYNC DATA')
    tb(s, 'One placement changes the trajectory of a career.',
       95, 160, W - 190, 110, 36, WHITE, bold=True, wrap=True)

    stats = [
        ('812%',         'Streams jump the week a track airs on a major TV show'),
        ('$5k – $500k',  'Sync fees per placement depending on project and territory'),
        ('$178M',        'US sync royalties in H1 2022 alone — up 30% year on year'),
    ]
    sy = 292
    for val, lbl in stats:
        add_rect(s, 95, sy, W - 190, 116, fill=STAT_BG, border=STAT_BD, bpt=0.75, radius_px=28)
        tb(s, val, 135, sy + 22, 300, 72, 64, LIME, bold=True)
        tb(s, lbl,  450, sy + 30, W - 570, 56, 25, DIM_P, wrap=True)
        sy += 132

    tb(s, 'Sources: Blakmarigold / industry data · Trolley sync royalty report',
       95, sy + 8, W - 190, 36, 18, SRC_CLR, italic=True)


def s09(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, PURPLE)
    # No glow on slide 9
    chrome(s, 9, 10, 'THE GAP')
    tb(s, 'African composers make the music the world wants.',
       95, 185, W - 190, 330, 84, WHITE, bold=True)
    tb(s, 'But most never see a sync deal.',
       95, 535, W - 190, 72, 38, WHITE)
    tb(s, 'No agent. No publisher. No pathway in.',
       95, 615, W - 190, 55, 30, DIM_P)
    tb(s, 'The music is world-class.',
       95, 702, W - 190, 72, 44, LIME, bold=True)
    tb(s, "The infrastructure isn't. Yet.",
       95, 780, W - 190, 72, 44, LIME, bold=True)


def s10(prs, blank):
    s = prs.slides.add_slide(blank)
    solid_bg(s, PURPLE)
    glow(s, 0.46, 640, 0.34)
    chrome(s, 10, 10, 'JOIN NOW', swipe=False, center_wm=True)

    tb(s, "That's the gap we're closing.",
       95, 295, W - 190, 210, 96, WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, ('SyncMaster connects African composers to global sync briefs '
           '— curated, vetted, rights-ready.'),
       95, 510, W - 190, 155, 34, SUB_10, align=PP_ALIGN.CENTER, wrap=True)

    # CTA lime button
    bw, bh = 530, 100
    bl = (W - bw) / 2
    add_rect(s, bl, 682, bw, bh, fill=LIME, radius_px=50)
    tb(s, 'Apply as a composer  →',
       bl, 710, bw, 52, 40, INK, bold=True, align=PP_ALIGN.CENTER)

    tb(s, 'Link in bio · Takes 60 seconds',
       95, 805, W - 190, 50, 26, CAP_10, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# Build
# ════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = e(W)
    prs.slide_height = e(H)
    blank = prs.slide_layouts[6]

    s01(prs, blank)
    s02(prs, blank)
    s03(prs, blank)
    s04(prs, blank)
    s05(prs, blank)
    s06(prs, blank)
    s07(prs, blank)
    s08(prs, blank)
    s09(prs, blank)
    s10(prs, blank)

    out = (r"C:\Users\infon\Downloads\Syncmaster-carousel-handoff-unzip"
           r"\syncmaster-carousel\SyncMaster_Carousel_02_Editable.pptx")
    prs.save(out)
    kb = os.path.getsize(out) // 1024
    print(f"Done — {len(prs.slides)} slides, {kb} KB")
    print(f"Path: {out}")


if __name__ == '__main__':
    build()
